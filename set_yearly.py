from functionUtily import *
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
sys.setrecursionlimit(10000)

def set_yearly (prs, mon_unm_list, yearNum, pageCon) : 
    
    targetPage = pageCon["targetPage"] -1
    repeat = pageCon["repeat"]
    source_slide_index = pageCon["source_slide_index"]
    link_page_num = pageCon["link_page_num"]
    if (source_slide_index == 0) :
        source_slide_index = targetPage
    for i in range(repeat):
        source_slide_index = source_slide_index + i
        source_slide = prs.slides[source_slide_index]
        source_slide_shapes = source_slide.shapes;  
        for shape in source_slide_shapes:  
            if (shape.name == "D" and MSO_SHAPE_TYPE.GROUP == shape.shape_type ):    
                1==1
                for month_shape in shape.shapes : 
                    month_shape_name = month_shape.name
                    monthinfo = int(month_shape_name)

                    cur_month_day_num = mon_unm_list[monthinfo-1];
                    eve_month_day_num = mon_unm_list[monthinfo-2];
                    month_start_week = fun_get_start_week_this_month(yearNum, monthinfo); # 1월 => 1                        
                    month_end_week = fun_get_end_week_this_month(yearNum, monthinfo, cur_month_day_num); #4
                    
                    start_day = 1;
                    next_day = 0;
                    cur_start_shape_num = month_start_week +1
                    eve_start_day = eve_month_day_num-cur_start_shape_num
                    if (month_end_week < 7) : 
                        next_day = 6 - month_end_week # 2
                        month_end_week = month_end_week + 1

                    last_day = next_day + 1 # 3                      
                    
                    count = 0;
                    for day_shape in month_shape.shapes :
                        cur_day = 0
                        if (month_start_week > 0) : 
                            cur_day = eve_start_day
                            if monthinfo == 1 : 
                                cur_day = 0;
                                count = 1
                            if count == 0 : 
                                link_page_num = link_page_num - 7
                                count = 1
                            eve_start_day = eve_start_day + 1
                            month_start_week = month_start_week -1
                        elif (month_start_week == 0 and cur_month_day_num > 0 ) : 
                            cur_day = start_day
                            start_day = start_day + 1
                            cur_month_day_num = cur_month_day_num -1
                            cur_start_shape_num = cur_start_shape_num + 1 
                        elif (cur_month_day_num == 0 and next_day > 0 and monthinfo < 12) : 
                            cur_day = last_day - next_day # 3 - 2
                            next_day = next_day -1
                                                    
                        if (cur_day == 0) :
                            link_page = targetPage
                        else :
                            link_page =link_page_num-1
                            link_page_num = link_page_num+1
                            
                        set_hyperlink_simple (day_shape, link_page, prs)
                            