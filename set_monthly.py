from pptx.enum.shapes import MSO_SHAPE_TYPE
from functionUtily import *
from pptx.enum.text import PP_ALIGN 
import sys
sys.setrecursionlimit(10000)


def set_monthly (prs, mon_unm_list, yearNum, pageCon, hoilydae_list): 
    
    targetPage = pageCon["targetPage"] -1
    repeat = pageCon["repeat"]
    source_slide_index = pageCon["source_slide_index"]
    link_page_num = pageCon["link_page_num"]
    sindex = 0;
    if (source_slide_index == 0) :
        sindex = targetPage
    else : sindex = source_slide_index
    monthinfo = 0;               
    count = 0
    for i in range(repeat):
        source_slide_index = sindex + i
        source_slide = prs.slides[source_slide_index]
        source_slide_shapes = source_slide.shapes;     
        for shape in source_slide_shapes:  
                
            1==1
            if (shape.name == "D" and MSO_SHAPE_TYPE.GROUP == shape.shape_type ):
                monthinfo = i + 1;
                cur_month_day_num = mon_unm_list[monthinfo-1];
                cur_month_last_day  = cur_month_day_num;
                eve_month_day_num = mon_unm_list[monthinfo-2];
                if (monthinfo == 1) : eve_month_day_num = mon_unm_list[monthinfo-1];
                
                month_start_week = fun_get_start_week_this_month(yearNum, monthinfo); # 1월 => 1                        
                month_end_week = fun_get_end_week_this_month(yearNum, monthinfo, cur_month_day_num);
                
                start_day = 1;
                next_day = 0;
                cur_start_shape_num = month_start_week +1
                eve_start_day = eve_month_day_num - month_start_week + 1
                if (month_end_week < 7) : 
                    next_day = 6 - month_end_week
                    month_end_week = month_end_week + 1

                last_day = next_day + 1 # 5         
                cur_day_txt = 0;
                # print(f" {monthinfo} 월 || eve_start_day : {eve_start_day}  || start_day : {start_day} || last_day : {last_day} || cur_month_day_num : {cur_month_day_num}")
                for day_shape in shape.shapes :  
                  
                    rGBColor = RGBColor(70, 70, 70)
                    cur_day = 0                 
                    if (month_start_week > 0) : 
                        cur_day = eve_start_day
                        last_month = monthinfo -1
                        last_year = yearNum
                        if monthinfo == 1 : 
                            cur_day = 0
                            count = 1
                            last_month = 12;
                            last_year = yearNum-1
                        if count == 0 : 
                            link_page_num = link_page_num - 5 #7
                            count = 1
                        cur_day_txt = eve_start_day
                        eve_start_day = eve_start_day + 1
                        month_start_week = month_start_week -1
                        
                        month_end_week = fun_get_day_color(last_year, last_month , cur_day_txt, hoilydae_list);
                        if (month_end_week == 0) : rGBColor = RGBColor(225, 155, 155)
                        else : rGBColor = RGBColor(166, 166, 166)
                    elif (month_start_week == 0 and cur_month_day_num > 0 ) : 
                        count = 2
                        cur_day = start_day
                        cur_day_txt = start_day
                        start_day = start_day + 1
                        cur_month_day_num = cur_month_day_num -1
                        cur_start_shape_num = cur_start_shape_num + 1 
                        if (cur_day_txt == cur_month_last_day) : count = 0
                        # print(f"link_page : {link_page}  || cur_day_txt : {cur_day_txt} || cur_day : {cur_day} || count : {count}")

                        month_end_week = fun_get_day_color(yearNum, monthinfo , cur_day_txt, hoilydae_list);
                        if (month_end_week == 0) : rGBColor = RGBColor(192, 0, 0)
                    elif (cur_month_day_num == 0 and next_day > 0) : 
                        count = 0
                        cur_day = last_day - next_day
                        # if (monthinfo == 2 and cur_day == 1) : link_page_num = link_page_num +1;
                        next_day = next_day -1
                        cur_day_txt = cur_day
                        next_month = monthinfo + 1
                        next_year = yearNum
                        if ( monthinfo == 12) : 
                            cur_day = 0
                            next_month = 1
                            next_year = yearNum+1
                        month_end_week = fun_get_day_color(next_year, next_month , cur_day_txt, hoilydae_list);
                        if (month_end_week == 0) : rGBColor = RGBColor(225, 155, 155)
                        else : rGBColor = RGBColor(166, 166, 166)
                                                
                    if (count == 2 and day_shape.name == "1") : 
                        link_page_num = link_page_num + 2
                        
                    if (cur_day == 0) :
                        link_page = targetPage
                    else :
                        link_page =link_page_num-1
                        link_page_num = link_page_num+1
                    
                                
                    src_text_frame = day_shape.text_frame    
                    fontSize = 0
                    fontName =""
                    for run in src_text_frame.paragraphs[0].runs:
                        fontSize = run.font.size
                        fontName = run.font.name
                        
                    # print(f"link_page : {link_page}  || cur_day_txt : {cur_day_txt} || cur_day : {cur_day} || count : {count}")
                    set_hyperlink_simple (day_shape, link_page, prs)
                    day_shape.text = str(cur_day_txt)
                    source_paragraph = day_shape.text_frame.paragraphs[0]
                    for run in source_paragraph.runs:
                        new_run_font = run.font
                        new_run_font.name = fontName
                        new_run_font.size = fontSize
                        new_run_font.color.rgb =rGBColor
                    source_paragraph.alignment = PP_ALIGN.CENTER
                link_page_num = link_page