from functionUtily import *
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
import sys
sys.setrecursionlimit(10000)

def set_weekly (prs, mon_unm_list, yearNum, pageCon, Month_Ko, hoilydae_list) : 
    targetPage = pageCon["targetPage"] -1
    repeat = pageCon["repeat"]

    month_link_page_num = pageCon["month_link_page_num"] -1

    source_slide_index = pageCon["source_slide_index"]
    link_page_num = pageCon["link_page_num"]
    if (source_slide_index == 0) :
        source_slide_index = targetPage

    source_slide_index = source_slide_index -1    
    source_slide = prs.slides[source_slide_index]
    source_slide_shapes = source_slide.shapes;    
    c = 0
    for shape in source_slide_shapes:  
        M_shapes = None
        1==1
        if (shape.name == "D" and MSO_SHAPE_TYPE.GROUP == shape.shape_type ):
            
            for source_slide_shape in source_slide_shapes :
                if (source_slide_shape.name == "M") :
                    M_shapes = source_slide_shape
            if ( M_shapes ) : 
                link_num_up = 0;
                monthinfo = 1; #1윌
                nextTimigW = 0;
                nextTimig = 4; #1일 
                cur_month_last_day = 2;
                hyperlink_page = 0;
                link_num = link_page_num-1
                eve_month_start_day = 31
                M_shapes_text = Month_Ko[0];
                for i in range(repeat):
                    target_Page_num = targetPage + i;
                    if (nextTimigW == nextTimig) : # ex. 1월 31일인데, dayinfo 가 32일이 되면 2월이 되고, dayinfo는 다시 1이 됨.
                        if(monthinfo == 12) : break;
                        M_shapes_text = Month_Ko[monthinfo]
                        month_link_page_num = month_link_page_num + 1
                        monthinfo = monthinfo + 1;
                        nextTimigW = 0;
                        link_num = link_page_num 
                        eve_month_start_day = cur_month_last_day;
                        c = 0
                    nextTimigW = nextTimigW + 1; 
                    link_page_num = link_num 
                    cur_month_day_num = mon_unm_list[monthinfo-1]; # 31  
                    cur_month_last_day = cur_month_day_num
                    month_start_week = fun_get_start_week_this_month(yearNum, monthinfo); # 시작하는 요일 (일 : 0, 월 : 1 ....)   
                      

                    month_end_week = fun_get_end_week_this_month(yearNum, monthinfo, cur_month_day_num);  # 이번달이 끝나는 요일 (일 : 0, 월 : 1 ....)  0 이면     
                    next_month_day_num = 6 - (month_end_week-1) # 일요일이면 6

                    
                    if (cur_month_last_day == 31 and month_start_week > 4) :
                        nextTimig = 5
                    elif(cur_month_last_day == 30 and month_start_week > 5) :
                        nextTimig = 5
                    elif(monthinfo == 12) :
                        nextTimig = 5
                    elif (monthinfo != 2  and month_end_week == 6) :
                        nextTimig = 5
                    else :
                        nextTimig = 4

                    target_slide = prs.slides[target_Page_num];
                    pg=target_Page_num                    
                            
                    if M_shapes.has_text_frame:
                        new_M_shape = copy_text_box (month_link_page_num, M_shapes, target_slide, prs);
                        new_M_shape.text = M_shapes_text;
                        new_text_frame = new_M_shape.text_frame
                        # 텍스트와 텍스트 스타일 복사
                        for i, paragraph in enumerate(M_shapes.text_frame.paragraphs):
                            if i >= len(new_text_frame.paragraphs):
                                new_text_frame.add_paragraph()
                            run = paragraph.runs[0];
                            for new_p in new_text_frame.paragraphs:
                                new_p.alignment = paragraph.alignment
                                for new_run in new_p.runs :
                                    new_run.font.size = run.font.size
                                    new_run.font.bold = run.font.bold
                                    new_run.font.italic = run.font.italic
                                    new_run.font.color.rgb = run.font.color.rgb
                                    new_run.font.name = run.font.name
                    day_txt = 0;
                    next_month_start_day = next_month_day_num +1 

                    if (month_start_week > 0 and monthinfo > 1) : link_page_num = link_page_num - 7   
                    for day_shape in shape.shapes :
                        link_num_up = 0;
                        rGBColor = RGBColor(70, 70, 70)
                        hyperlink_page = link_page_num;
                        # 이번달
                        if (month_start_week == 0 and cur_month_day_num > 0) :
                            link_num_up = 1;
                            cur_month_day_num = cur_month_day_num -1; # 30
                            day_txt = cur_month_last_day - cur_month_day_num; # 31-30 1일
                            month_end_week = fun_get_day_color(yearNum, monthinfo , day_txt, hoilydae_list);
                            if (month_end_week == 0) : rGBColor = RGBColor(192, 0, 0)
                        # 전달/다음달
                        else : 
                            #전달
                            if (month_start_week > 0) :
                                link_num_up = 1; 
                                day_txt = eve_month_start_day - (month_start_week-1)
                                if (monthinfo == 1) : 
                                    hyperlink_page = pg;
                                    # print(f" {yearNum} 년 {monthinfo} 월 {day_txt} 일")       
                                    month_end_week = fun_get_day_color(yearNum-1, 12 , day_txt, hoilydae_list);
                                    link_num_up = 0;
                                else : month_end_week = fun_get_day_color(yearNum, monthinfo-1 , day_txt, hoilydae_list);
                                month_start_week = month_start_week -1;
                            # 다음 달
                            elif (cur_month_day_num == 0 and next_month_day_num > 0) : 
                                link_num_up = 1; 
                                day_txt = next_month_start_day - (next_month_day_num);
                                if (monthinfo == 12) : 
                                    hyperlink_page = pg;
                                    month_end_week = fun_get_day_color(yearNum+1, 1 , day_txt, hoilydae_list);
                                else : month_end_week = fun_get_day_color(yearNum, monthinfo+1 , day_txt, hoilydae_list);
                                next_month_day_num = next_month_day_num - 1
                            if (month_end_week == 0) : rGBColor = RGBColor(225, 155, 155)
                            else : rGBColor = RGBColor(166, 166, 166)
                        new_txt = str(day_txt)
                        if (month_start_week == 0  and cur_month_day_num == 0 and next_month_day_num == 0) : 
                            new_txt = ""
                            hyperlink_page = pg
                            link_num_up = 0; 
                        new_shape = copy_text_box(hyperlink_page, day_shape, target_slide, prs);

                        if day_shape.has_text_frame:
                            new_shape.text = new_txt;

                            # 텍스트 포맷 복사
                            for runs, new_runs in zip(day_shape.text_frame.paragraphs, new_shape.text_frame.paragraphs):
                                for p, new_p in zip (runs.runs, new_runs.runs ) : 
                                    new_p.font.size = p.font.size
                                    new_p.font.bold = p.font.bold
                                    new_p.font.name = p.font.name
                                    new_p.font.italic = p.font.italic
                                    new_p.font.color.rgb = rGBColor
                                new_runs.alignment = runs.alignment
                        
                            new_shape.shadow.inherit = False
                            new_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                        link_page_num = link_page_num + link_num_up;    
                        # print(f" {monthinfo} 월 || pg:{pg} ||day_txt : {new_txt} || hyperlink_page : {hyperlink_page}|| week_link_page_num : {week_link_page_num} || link_page_num : {link_page_num}")  
                    