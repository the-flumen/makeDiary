from functionUtily import *
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
import sys
sys.setrecursionlimit(10000)

def set_weekly_money (prs, mon_unm_list, yearNum, pageCon, hoilydae_list) : 
    targetPage = pageCon["targetPage"] -2
    repeat = pageCon["repeat"]
    source_slide_index = pageCon["source_slide_index"]
    link_page_num = pageCon["link_page_num"] -2
    if (source_slide_index == 0) :
        source_slide_index = targetPage
    monthinfo = 0;


    eve_year_repeat_num = fun_get_start_week_this_month(yearNum, 1);
    month_1_Start_week = eve_year_repeat_num - 1 ; # 월요일 1
    link_num_up_is = False;
    cur_month_day_num = 31 
    day_text = 31 - month_1_Start_week;  
    if (month_1_Start_week < 0) : 
        day_text = 1
        yearNum = yearNum +1;
    for i in range(repeat):
        source_slide_index = source_slide_index + 1
        source_slide = prs.slides[source_slide_index]
        source_slide_shapes = source_slide.shapes;
        week_shapes = None
        week_ko_shapes = None   

        for shape in source_slide_shapes:  
            if (shape.name == "W" and MSO_SHAPE_TYPE.GROUP == shape.shape_type ):
                week_shapes = shape
            if (shape.name == "WK" and MSO_SHAPE_TYPE.GROUP == shape.shape_type ):
                week_ko_shapes = shape

        if week_shapes and week_ko_shapes:
            w = 0;  
            w = w + 1; # 1주         
            for W_shape, week_ko_shape in zip(week_shapes.shapes, week_ko_shapes.shapes):
                rGBColor = RGBColor(70, 70, 70);
                cur_week = 0

                if (eve_year_repeat_num > 0) : 
                    
                    eve_year_repeat_num = eve_year_repeat_num - 1
                    cur_week = fun_get_day_color(yearNum-1, 12, day_text, hoilydae_list); 
                else :
                    link_num_up_is = True;
                    cur_month_day_num = mon_unm_list[monthinfo-1];
                    cur_week = fun_get_day_color(yearNum, monthinfo , day_text, hoilydae_list);
                if (cur_week == 0) : rGBColor = RGBColor(192, 0, 0)
                src_text_frame = W_shape.text_frame    
                fontSize = 0
                fontName =""
                for run in src_text_frame.paragraphs[0].runs:
                    fontSize = run.font.size
                    fontName = run.font.name
                    
                    
                if (link_num_up_is) : 
                    link_page_num = link_page_num + 1;
                    set_hyperlink_simple (W_shape, link_page_num, prs)
                else :   
                    set_hyperlink_simple (W_shape, source_slide_index, prs)
                
                W_shape.text = str(day_text)
                source_paragraph = W_shape.text_frame.paragraphs[0]
                source_ko_paragraph = week_ko_shape.text_frame.paragraphs[0]
                
                for run, k_run in zip(source_paragraph.runs, source_ko_paragraph.runs):
                    new_run_font = run.font
                    new_run_font.name = fontName
                    new_run_font.size = fontSize
                    new_run_font.color.rgb =rGBColor
                    k_run.font.color.rgb =rGBColor
                source_paragraph.alignment = PP_ALIGN.CENTER
                if(cur_month_day_num == day_text) : 
                    if (monthinfo == 12) : 
                        link_num_up_is = False;
                        monthinfo = 1
                        yearNum = yearNum + 1
                    day_text = 0
                    monthinfo = monthinfo + 1; #달
                day_text = day_text + 1;                