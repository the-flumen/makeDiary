
import datetime
from datetime import datetime
from pytimekr import pytimekr
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor

def set_textbox (source_paragraph, new_paragraph, RGBColor):
    for run in source_paragraph.runs:
        for new_run in new_paragraph.runs:
            new_run_font = new_run.font
            target_font = run.font
            new_run_font.name = target_font.name
            new_run_font.size = target_font.size
            if target_font.size == 838200 : 
                new_run_font.size = 889000
            if target_font.size == (8*12700) :
                new_run_font.size = (9*12700)
            new_run_font.color.rgb =RGBColor
            new_run_font.spacing = Pt(1)
    new_paragraph.alignment = PP_ALIGN.CENTER

def fun_get_end_week_this_month (yearNum, monthinfo, dayinfo) :
    start_day = datetime(yearNum, monthinfo, dayinfo)
    weekday_number = start_day.weekday(); 
    
    if (weekday_number == 6) :
        return 0
    #"Return day of the week, where Monday == 0 ... Sunday == 6."
    return weekday_number+1; #월요일이면 0+2, 화요일 1 => 3

def fun_get_day_color (yearNum, monthinfo, dayinfo, hoilydae_list) :
    start_day = datetime(yearNum, monthinfo, dayinfo)
    weekday_number = start_day.weekday(); 
    
    if (weekday_number == 6) :
        return 0
    
    holiday_li = pytimekr.holidays(year = yearNum)
    for i in holiday_li:
        if(start_day.strftime("%Y-%m-%d") == str(i)): 
            return 0
    for i in hoilydae_list:
        if(start_day.strftime("%Y-%m-%d") == str(i)): 
            return 0
    return weekday_number+1; #월요일이면 0+2, 화요일 1 => 3

def fun_is_hoilday (start_day, holiday_li) :
    for i in holiday_li:
        if(start_day.strftime("%Y-%m-%d") == str(i)): 
            return 0

def fun_get_start_week_this_month (yearNum, monthinfo) :
    start_day = datetime(yearNum, monthinfo, 1)
    weekday_number = start_day.weekday(); 
    if (weekday_number == 6) :
        return 0
    
    #"Return day of the week, where Monday == 0 ... Sunday == 6."
    return weekday_number+1; #월요일이면 0+1, 화요일 1+1 => 3


def fun_get_strat_repeat_num (yearNum) :
    # days_of_week = ['월요일', '화요일', '수요일', '목요일', '금요일', '토요일', '일요일']
    start_day = datetime(yearNum, 1, 1)
    weekday_number = start_day.weekday(); 
    print(weekday_number);
    if (weekday_number == 6) :
        return 7
    # 0: 월 ~ 6: 1
    return 6 - weekday_number; #월요일이면 6

def set_hyperlink_simple (source_shape, link_page_num, prs):
    link_page_slide = prs.slides[link_page_num]
    copy_hyperlink(link_page_slide, source_shape.click_action)   
    
def set_hyperlink (source_shape, hyperlink_page, i, prs):
    new_text_frame = source_shape.text_frame    
    fontSize = 0
    for run in new_text_frame.paragraphs[0].runs:
        fontSize = run.font.size
    r = 0
    slide2 = prs.slides[1]
    if  fontSize == 889000 :        
        slide2 = prs.slides[hyperlink_page]
    elif fontSize < 127000:
        index_num = (2*(int(source_shape.text)-1) + i)
        slide2 = prs.slides[index_num]
        if (r < index_num) :
            r = index_num
    
    copy_hyperlink(slide2, source_shape.click_action)

def copy_hyperlink (source_hyperlink, target_click_action): 
    target_click_action.target_slide = source_hyperlink


def copy_shpes_inc (hyperlink_page, source_shape, target_slide, prs):          
    left = source_shape.left / Inches(1)
    top = source_shape.top / Inches(1)
    width = source_shape.width / Inches(1)
    height = source_shape.height / Inches(1)
    position_adjustment_cm = 9.14
    cm_to_inches = 2.54
    position_adjustment_in = position_adjustment_cm / cm_to_inches
    new_shape = target_slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(left - position_adjustment_in),  # 조정된 왼쪽 위치
            Inches(top),  # 원본 도형과 동일한 위쪽 위치
            Inches(width),  # 원본 도형과 동일한 너비
            Inches(height)  # 원본 도형과 동일한 높이
        )
    new_shape.line.color.rgb = RGBColor(0, 0, 0)
    new_shape.line.fill.background()
    new_shape.fill.background()  
    if (hyperlink_page != -1 ):
        hyperlink_page_slide = prs.slides[hyperlink_page] 
        new_shape.click_action.target_slide = hyperlink_page_slide
    return new_shape


def copy_shpes (hyperlink_page, source_shape, target_slide, prs):          
    left = source_shape.left
    top = source_shape.top
    width = source_shape.width
    height = source_shape.height
    new_shape = target_slide.shapes.add_shape(MSO_SHAPE.OVAL,
            left,
            top,
            width,
            height
        )
    new_shape.line.color.rgb = RGBColor(0, 0, 0)
    new_shape.line.fill.background()
    new_shape.fill.background()  
    if (hyperlink_page != -1 ):
        hyperlink_page_slide = prs.slides[hyperlink_page] 
        new_shape.click_action.target_slide = hyperlink_page_slide
    return new_shape

# def copy_shpes (hyperlink_page, source_shape, target_slide, prs):
#     # copy_shpes(target_slide_index, source_shape, target_slide, prs)         
#     return create_new_shpes (source_shape, target_slide, hyperlink_page, prs)



def set_text_box (source_shape, new_shape, target_slide, new_txt):  

    new_shape.text = new_txt;

    # 텍스트 포맷 복사
    for runs, new_runs in zip(source_shape.text_frame.paragraphs, new_shape.text_frame.paragraphs):
        for p, new_p in zip (runs.runs, new_runs.runs ) : 
            new_p.font.size = p.font.size
            new_p.font.bold = p.font.bold
            new_p.font.name = p.font.name
            new_p.font.italic = p.font.italic
            new_p.font.color.rgb = p.font.color.rgb
        new_runs.alignment = runs.alignment

    new_shape.shadow.inherit = False
    new_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def copy_text_box (hyperlink_page, source_shape, target_slide, prs):          
    left = source_shape.left
    top = source_shape.top
    width = source_shape.width
    height = source_shape.height
    new_shape = target_slide.shapes.add_textbox(
            left,
            top,
            width,
            height
        )
    new_shape.line.color.rgb = RGBColor(0, 0, 0)
    new_shape.line.fill.background()
    new_shape.fill.background()  
    if (hyperlink_page != -1 ):
        hyperlink_page_slide = prs.slides[hyperlink_page] 
        new_shape.click_action.target_slide = hyperlink_page_slide
    return new_shape
