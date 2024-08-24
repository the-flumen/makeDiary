from pptx import Presentation
from functionUtily import *
import sys
sys.setrecursionlimit(10000)


def fun_get_strat_repeat_num (yearNum) :
    # days_of_week = ['월요일', '화요일', '수요일', '목요일', '금요일', '토요일', '일요일']
    start_day = datetime(yearNum, 1, 1)
    weekday_number = start_day.weekday(); 
    print(weekday_number);
    if (weekday_number == 6) :
        return 7
    # 0: 월 ~ 6: 1
    return 6 - weekday_number; #월요일이면 6

def extract_W(yearNum, keyword, input_file, output_file, page_config, start_slide_index, start_link_slid_index):
    try:
        prs = Presentation(input_file)
        source_slide = prs.slides[start_slide_index]
        
        for shape in source_slide.shapes:
        
            if (keyword == shape.name) : 
                shapes_to_copy = shape
                # 슬라이드의 도형과 텍스트 추출
                for pageCon in page_config :
                    cunt = 1;
                    startPageCon = pageCon["startPage"]
                    repeat_num = fun_get_strat_repeat_num (yearNum);
                    hyperlink_page = start_link_slid_index
                    for i in range(366):
                        if (repeat_num == 0) : 
                            repeat_num = 7
                            hyperlink_page = start_link_slid_index + cunt;
                            cunt = cunt+1;
                        target_slide_index = startPageCon + i;
                        # print(f" i : {i} || hyperlink_page : {hyperlink_page}  || target_slide_index : {target_slide_index} || repeat_num : {repeat_num}")
                        if (target_slide_index != start_slide_index) :
                            target_slide = prs.slides[target_slide_index];
                            copy_shpes(hyperlink_page, shapes_to_copy, target_slide, prs)
                        repeat_num = repeat_num - 1;
                            
        prs.save(output_file)
        print()
        print(f"파일 '{output_file}'이(가) 생성되었습니다.")

    except Exception as e:
        print(f"파워포인트 파일을 열거나 데이터를 추출하는 중 오류가 발생했습니다: {e}")
        return None
    



